import argparse
import csv
import html
import importlib.util
import json
import pathlib
import re
import sys


def has_module(name):
    return importlib.util.find_spec(name) is not None


def probe():
    engines = {
        "docling": has_module("docling"),
        "markitdown": has_module("markitdown"),
        "pypdf": has_module("pypdf"),
        "python_docx": has_module("docx"),
        "openpyxl": has_module("openpyxl"),
        "python_pptx": has_module("pptx"),
    }
    print(json.dumps({
        "python": sys.executable,
        "engines": engines,
        "ready": any(engines.values()),
        "message": "Docling preferred, MarkItDown fallback, then lightweight local extractors."
    }))


def parse_with_docling(source):
    from docling.document_converter import DocumentConverter

    result = DocumentConverter().convert(str(source))
    document = result.document
    markdown = document.export_to_markdown()
    structured = {
        "title": source.name,
        "type": source.suffix.lower().lstrip("."),
        "sections": [{"kind": "document", "heading": source.name, "text": markdown}],
        "docling": json.loads(document.export_to_json()) if hasattr(document, "export_to_json") else None,
    }
    return emit(source, "docling", markdown, structured)


def parse_with_markitdown(source):
    from markitdown import MarkItDown

    result = MarkItDown().convert(str(source))
    markdown = result.text_content
    structured = {
        "title": source.name,
        "type": source.suffix.lower().lstrip("."),
        "sections": [{"kind": "document", "heading": source.name, "text": markdown}],
    }
    return emit(source, "markitdown", markdown, structured)


def parse_pdf(source):
    from pypdf import PdfReader

    reader = PdfReader(str(source))
    sections = []
    lines = [f"# {source.name}", ""]
    for index, page in enumerate(reader.pages, start=1):
        text = (page.extract_text() or "").strip()
        if not text:
            continue
        heading = f"Page {index}"
        lines.extend([f"## {heading}", "", text, ""])
        sections.append({"kind": "page", "page": index, "heading": heading, "text": text})
    markdown = "\n".join(lines).strip()
    return emit(source, "pypdf", markdown, sections_to_structured(source, sections))


def parse_docx(source):
    import docx

    document = docx.Document(str(source))
    sections = []
    lines = [f"# {source.name}", ""]
    current_heading = source.name
    buffer = []
    for paragraph in document.paragraphs:
        text = paragraph.text.strip()
        if not text:
            continue
        style = (paragraph.style.name or "").lower()
        if "heading" in style:
            if buffer:
                sections.append({"kind": "section", "heading": current_heading, "text": "\n\n".join(buffer)})
                buffer = []
            current_heading = text
            lines.extend([f"## {text}", ""])
        else:
            buffer.append(text)
            lines.extend([text, ""])
    if buffer:
        sections.append({"kind": "section", "heading": current_heading, "text": "\n\n".join(buffer)})

    for table_index, table in enumerate(document.tables, start=1):
        table_md = table_to_markdown([[cell.text.strip() for cell in row.cells] for row in table.rows])
        if table_md:
            heading = f"Table {table_index}"
            lines.extend([f"## {heading}", "", table_md, ""])
            sections.append({"kind": "table", "heading": heading, "text": table_md})
    return emit(source, "python-docx", "\n".join(lines).strip(), sections_to_structured(source, sections))


def parse_xlsx(source):
    import openpyxl

    workbook = openpyxl.load_workbook(str(source), read_only=True, data_only=True)
    sections = []
    lines = [f"# {source.name}", ""]
    for sheet in workbook.worksheets:
        rows = []
        for row in sheet.iter_rows(values_only=True):
            values = ["" if value is None else str(value) for value in row]
            if any(value.strip() for value in values):
                rows.append(values)
            if len(rows) >= 500:
                break
        table = table_to_markdown(rows)
        if table:
            lines.extend([f"## {sheet.title}", "", table, ""])
            sections.append({"kind": "sheet", "heading": sheet.title, "text": table})
    return emit(source, "openpyxl", "\n".join(lines).strip(), sections_to_structured(source, sections))


def parse_pptx(source):
    from pptx import Presentation

    deck = Presentation(str(source))
    sections = []
    lines = [f"# {source.name}", ""]
    for index, slide in enumerate(deck.slides, start=1):
        texts = []
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                texts.append(shape.text.strip())
        if texts:
            text = "\n\n".join(texts)
            heading = f"Slide {index}"
            lines.extend([f"## {heading}", "", text, ""])
            sections.append({"kind": "slide", "slide": index, "heading": heading, "text": text})
    return emit(source, "python-pptx", "\n".join(lines).strip(), sections_to_structured(source, sections))


def parse_image(source):
    raise RuntimeError("Image OCR requires Docling or MarkItDown with OCR dependencies installed.")


def sections_to_structured(source, sections):
    return {
        "title": source.name,
        "type": source.suffix.lower().lstrip("."),
        "sections": sections,
    }


def table_to_markdown(rows):
    rows = [row for row in rows if any(str(cell).strip() for cell in row)]
    if not rows:
        return ""
    width = max(len(row) for row in rows)
    normalized = [[html.escape(str(cell)).replace("\n", " ").strip() for cell in row] + [""] * (width - len(row)) for row in rows]
    header = normalized[0]
    separator = ["---"] * width
    body = normalized[1:]
    return "\n".join([
        "| " + " | ".join(header) + " |",
        "| " + " | ".join(separator) + " |",
        *["| " + " | ".join(row) + " |" for row in body],
    ])


def emit(source, parser, markdown, structured):
    markdown = clean_markdown(markdown)
    print(json.dumps({
        "parser": parser,
        "title": source.name,
        "type": source.suffix.lower().lstrip("."),
        "markdown": markdown,
        "structured": structured,
    }, ensure_ascii=False))


def clean_markdown(markdown):
    text = markdown.replace("\r\n", "\n")
    text = re.sub(r"\n{4,}", "\n\n\n", text)
    return text.strip()


def parse(source):
    suffix = source.suffix.lower()
    if has_module("docling"):
        try:
            return parse_with_docling(source)
        except Exception as exc:
            print(f"Docling failed, falling back: {exc}", file=sys.stderr)
    if has_module("markitdown"):
        try:
            return parse_with_markitdown(source)
        except Exception as exc:
            print(f"MarkItDown failed, falling back: {exc}", file=sys.stderr)

    if suffix == ".pdf" and has_module("pypdf"):
        return parse_pdf(source)
    if suffix == ".docx" and has_module("docx"):
        return parse_docx(source)
    if suffix == ".xlsx" and has_module("openpyxl"):
        return parse_xlsx(source)
    if suffix == ".pptx" and has_module("pptx"):
        return parse_pptx(source)
    if suffix in {".png", ".jpg", ".jpeg", ".tiff", ".bmp", ".webp"}:
        return parse_image(source)
    raise RuntimeError(f"No parser available for {suffix}. Install docling or markitdown[all].")


def main():
    arg_parser = argparse.ArgumentParser()
    arg_parser.add_argument("source", nargs="?")
    arg_parser.add_argument("--probe", action="store_true")
    args = arg_parser.parse_args()
    if args.probe:
        return probe()
    if not args.source:
        raise RuntimeError("Missing source path.")
    parse(pathlib.Path(args.source))


if __name__ == "__main__":
    try:
        main()
    except Exception as exc:
        print(str(exc), file=sys.stderr)
        sys.exit(1)
